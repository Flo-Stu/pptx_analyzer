import os
import io
import requests
from flask import Flask, request, render_template, send_file, redirect, url_for
from pptx import Presentation

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'

# Stelle sicher, dass das Upload-Verzeichnis existiert
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/analyze', methods=['POST'])
def analyze():
    potx_url = request.form['potx_url']
    try:
        response = requests.get(potx_url)
        response.raise_for_status()
        potx_content = io.BytesIO(response.content)
        prs = Presentation(potx_content)
        layouts_info = extract_layouts_info(prs)
        return render_template('analyze.html', layouts=layouts_info, potx_url=potx_url)
    except requests.exceptions.RequestException as e:
        return f"Fehler beim Herunterladen der Datei: {e}"

def extract_layouts_info(prs):
    layouts_info = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [{'idx': ph.placeholder_format.idx, 'name': ph.name} for ph in layout.placeholders]
        layouts_info.append({'index': i, 'name': layout.name, 'placeholders': placeholders})
    return layouts_info

@app.route('/generate_script', methods=['POST'])
def generate_script():
    potx_url = request.form['potx_url']
    layouts_info = request.form.get('layouts_info')
    script_content = generate_python_script(potx_url, layouts_info)
    script_filename = 'generated_script.py'
    script_path = os.path.join(app.config['UPLOAD_FOLDER'], script_filename)
    with open(script_path, 'w', encoding='utf-8') as script_file:
        script_file.write(script_content)
    return send_file(script_path, as_attachment=True)

def generate_python_script(potx_url, layouts_info):
    script_lines = [
        "import io",
        "import requests",
        "from pptx import Presentation",
        "",
        "# Herunterladen der POTX-Datei",
        f"potx_url = '{potx_url}'",
        "response = requests.get(potx_url)",
        "response.raise_for_status()",
        "potx_content = io.BytesIO(response.content)",
        "prs = Presentation(potx_content)",
        "",
        "# Hinzufügen von Folien basierend auf den Layouts",
    ]
    for layout in layouts_info:
        script_lines.append(f"# Layout: {layout['name']}")
        script_lines.append(f"slide_layout = prs.slide_layouts[{layout['index']}]")
        script_lines.append("slide = prs.slides.add_slide(slide_layout)")
        for ph in layout['placeholders']:
            script_lines.append(f"# Platzhalter {ph['idx']}: {ph['name']}")
        script_lines.append("")
    script_lines.append("# Speichern der Präsentation")
    script_lines.append("prs.save('neue_praesentation.pptx')")
    return "\n".join(script_lines)

if __name__ == '__main__':
    app.run(debug=True)
